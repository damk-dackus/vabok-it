"""
Script om PowerPoint-presentaties te genereren voor de 5 workshops
van de Afstemmen en Plannen opleiding.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

OUTPUT_DIR = Path(__file__).parent / "presentaties"


def add_title_slide(prs, title, subtitle=""):
    """Voeg een titelslide toe."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle
    return slide


def add_content_slide(prs, title, bullet_points):
    """Voeg een inhoudsslide toe met bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {point}" if not point.startswith("•") else point
        p.level = 0
    return slide


def add_section_slide(prs, title, content_paragraphs):
    """Voeg een sectieslide toe met meerdere paragrafen."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, para in enumerate(content_paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = para
        p.level = 0
    return slide


def create_workshop_1():
    """Sessie 1: Kennismaking en wensen/eisen verzamelen"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Sessie 1: Kennismaking en wensen/eisen verzamelen",
        "Afstemmen en Plannen - Module Software Ontwikkeling"
    )

    add_content_slide(prs, "Les onderdeel - Overzicht", [
        "Kennismaking met elkaar",
        "Wat is projectopstart?",
        "Wat zijn wensen en eisen?",
        "Hoe verzamel je wensen en eisen?",
        "Wat zijn technische uitgangspunten?",
        "Wat is een programma van eisen?",
        "Overzicht van de opleiding"
    ])

    add_content_slide(prs, "Wat is projectopstart?", [
        "Het starten van een softwareproject met een duidelijke basis",
        "Doel: begrijpen wat er gebouwd moet worden",
        "Belangrijke eerste stap voordat je gaat ontwikkelen",
        "Zonder goede opstart loopt een project snel vast"
    ])

    add_content_slide(prs, "Wensen vs. eisen", [
        "Eisen: wat MOET het systeem kunnen (verplicht)",
        "Wensen: wat ZOU MOOI zijn om te hebben (optioneel)",
        "Functionele eisen: wat moet de software kunnen?",
        "Niet-functionele eisen: hoe moet de software werken? (snelheid, veiligheid, etc.)"
    ])

    add_content_slide(prs, "Hoe verzamel je wensen en eisen?", [
        "Spreek met de opdrachtgever en stakeholders",
        "Vraag door als iets niet duidelijk is",
        "Noteer alles zorgvuldig",
        "Gebruik het programma van eisen als documentatie"
    ])

    add_content_slide(prs, "Technische uitgangspunten", [
        "Technieken: programmeertaal, framework, database",
        "Code conventies",
        "Tools die je gaat gebruiken",
        "Documenteer alle keuzes en de motivatie ervoor"
    ])

    add_content_slide(prs, "Programma van eisen", [
        "Projectbeschrijving en doel",
        "Uitgangspunten",
        "Technische uitgangspunten",
        "Functionele eisen",
        "Niet-functionele eisen",
        "Wensen (nice-to-have)"
    ])

    add_content_slide(prs, "Opdracht deze sessie", [
        "Verzamel wensen en eisen voor je project",
        "Maak een programma van eisen (2-3 pagina's)",
        "Stem af met de opdrachtgever"
    ])

    add_content_slide(prs, "Online cursus", [
        "Project Management Foundations",
        "LinkedIn Learning",
        "Basis van projectmanagement en projectopstart"
    ])

    add_title_slide(prs, "Vragen?", "Tot de volgende sessie!")

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_DIR / "Sessie_1_Kennismaking_wensen_en_eisen.pptx")


def create_workshop_2():
    """Sessie 2: User stories en backlog maken"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Sessie 2: User stories en backlog maken",
        "Afstemmen en Plannen - Module Software Ontwikkeling"
    )

    add_content_slide(prs, "Les onderdeel - Overzicht", [
        "Wat zijn user stories?",
        "Hoe formuleer je een goede user story?",
        "Wat is een backlog?",
        "Wat is een definition of done?",
        "Hoe maak je een backlog?",
        "Voorbeelden van goede user stories"
    ])

    add_content_slide(prs, "Wat zijn user stories?", [
        "Korte beschrijving van een functionaliteit vanuit gebruikersperspectief",
        "Format: 'Als [rol] wil ik [functionaliteit] zodat [voordeel]'",
        "Vertaling van eisen naar implementeerbare items",
        "Basis voor de backlog en planning"
    ])

    add_content_slide(prs, "Een goede user story bevat", [
        "Wie: welke gebruiker/rol?",
        "Wat: welke functionaliteit?",
        "Waarom: welk voordeel of doel?",
        "Realistisch en testbaar"
    ])

    add_content_slide(prs, "Backlog en Definition of Done", [
        "Backlog: geordende lijst van alle user stories",
        "Definition of Done: wanneer is een user story 'klaar'?",
        "DoD zorgt voor kwaliteitsstandaarden",
        "Afstemmen met stakeholders is belangrijk"
    ])

    add_content_slide(prs, "PvE naar SRS", [
        "Programma van eisen → Software Requirement Specification",
        "PvE is business-gericht, SRS is technisch",
        "Functionele eisen (FE) en niet-functionele eisen (NFE)",
        "Meetbaar en testbaar formuleren"
    ])

    add_content_slide(prs, "Opdracht deze sessie", [
        "Zet programma van eisen om naar SRS",
        "Minimaal 10 functionele eisen (genummerd)",
        "Minimaal 5 niet-functionele eisen",
        "Kort reflectieverslag (1 pagina)"
    ])

    add_content_slide(prs, "Online cursus", [
        "Applying Agile MoSCoW Prioritization",
        "LinkedIn Learning",
        "MoSCoW: Must have, Should have, Could have, Won't have"
    ])

    add_title_slide(prs, "Vragen?", "Tot de volgende sessie!")

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_DIR / "Sessie_2_User_stories_en_backlog.pptx")


def create_workshop_3():
    """Sessie 3: Planning maken met MoSCoW"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Sessie 3: Planning maken met MoSCoW",
        "Afstemmen en Plannen - Module Software Ontwikkeling"
    )

    add_content_slide(prs, "Les onderdeel - Overzicht", [
        "Wat is prioriteren?",
        "Wat is de MoSCoW methode?",
        "Hoe gebruik je MoSCoW?",
        "Hoe maak je een planning?",
        "Hoe schat je tijd in?",
        "Hoe stem je een planning af met stakeholders?"
    ])

    add_content_slide(prs, "Prioriteren", [
        "Bepalen wat eerst moet gebeuren",
        "Niet alles tegelijk – focus op wat het belangrijkst is",
        "Helpt bij realistische planning en scope"
    ])

    add_content_slide(prs, "MoSCoW methode", [
        "Must have: Moet absoluut – zonder dit is het project geen succes",
        "Should have: Zou moeten – belangrijk maar niet kritisch",
        "Could have: Zou kunnen – nice-to-have als er tijd is",
        "Won't have: Gaat nu niet – bewust uit scope voor nu"
    ])

    add_content_slide(prs, "Planning maken", [
        "Must have user stories eerst in de planning",
        "Daarna Should have, dan Could have",
        "Houd rekening met afhankelijkheden",
        "Maak de planning realistisch"
    ])

    add_content_slide(prs, "Tijd inschatten", [
        "Per user story: hoeveel tijd kost het?",
        "Denk aan uren, dagen of weken",
        "Houd rekening met complexiteit",
        "Voeg buffer toe voor onverwachte zaken"
    ])

    add_content_slide(prs, "Afstemmen met stakeholders", [
        "Laat planning zien aan opdrachtgever",
        "Vraag of prioriteiten kloppen",
        "Vraag of planning realistisch is",
        "Pas aan waar nodig"
    ])

    add_content_slide(prs, "Opdracht deze sessie", [
        "Prioriteer user stories met MoSCoW",
        "Schat tijd in per user story",
        "Maak een planning (2-3 pagina's)",
        "Stem af met stakeholders"
    ])

    add_content_slide(prs, "Online cursus", [
        "Project Planning: Putting It All Together",
        "LinkedIn Learning",
        "Complete planning, tijd inschatten, afhankelijkheden"
    ])

    add_title_slide(prs, "Vragen?", "Tot de volgende sessie!")

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_DIR / "Sessie_3_Planning_MoSCoW.pptx")


def create_workshop_4():
    """Sessie 4: Ontwerpen met schema's"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Sessie 4: Ontwerpen met schema's",
        "Afstemmen en Plannen - Module Software Ontwikkeling"
    )

    add_content_slide(prs, "Les onderdeel - Overzicht", [
        "Wat is een ontwerp?",
        "Welke schema-technieken zijn er?",
        "Data design: ERD, class diagram, normalisatie",
        "User interface design: use case, wireframes, mock-ups",
        "Program logic design: activity diagram, flowchart",
        "AVG (GDPR) en OWASP top 10"
    ])

    add_content_slide(prs, "Data design", [
        "ERD (Entity Relationship Diagram): datamodel",
        "Class diagram: objectgeoriënteerd ontwerp",
        "Normalisatie: database optimalisatie",
        "Vraag: welke tabellen/klassen en relaties?"
    ])

    add_content_slide(prs, "User interface design", [
        "Use case diagram: welke acties kan de gebruiker doen?",
        "Wireframes: schets van schermen",
        "Mock-ups: gedetailleerd visueel ontwerp",
        "Vraag: welke schermen en navigatie?"
    ])

    add_content_slide(prs, "Program logic design", [
        "Activity diagram: procesflow",
        "Flowchart: logische stappen",
        "Vraag: welke stappen in welke volgorde?"
    ])

    add_content_slide(prs, "AVG (GDPR) – Privacy", [
        "Toestemming vragen voor gegevensverzameling",
        "Gegevens beveiligen",
        "Gebruikersrechten: inzage, correctie, verwijdering",
        "Ontwerp moet AVG-compliant zijn"
    ])

    add_content_slide(prs, "OWASP Top 10 – Veiligheid", [
        "Injectie aanvallen (SQL, etc.)",
        "Authenticatie problemen",
        "Gevoelige data blootstelling",
        "Ontwerp moet veiligheidsrisico's adresseren"
    ])

    add_content_slide(prs, "Opdracht deze sessie", [
        "Minimaal 3 schema's: data, UI, program logic",
        "Onderbouw ontwerpkeuzes",
        "Beschrijf AVG en OWASP compliance",
        "Ontwerp document (min. 3 pagina's)"
    ])

    add_content_slide(prs, "Online cursus", [
        "UML Foundations",
        "LinkedIn Learning",
        "Basis van UML en diagrammen"
    ])

    add_title_slide(prs, "Vragen?", "Tot de volgende sessie!")

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_DIR / "Sessie_4_Ontwerpen_met_schemas.pptx")


def create_workshop_5():
    """Sessie 5: Voortgang bewaken en evalueren"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Sessie 5: Voortgang bewaken en evalueren",
        "Afstemmen en Plannen - Module Software Ontwikkeling"
    )

    add_content_slide(prs, "Les onderdeel - Overzicht", [
        "Waarom is voortgang bewaken belangrijk?",
        "Hoe monitor je voortgang?",
        "Hoe controleer je of alles volgens planning verloopt?",
        "Hoe bewaak je prioriteiten?",
        "Hoe documenteer je afwijkingen?",
        "Hoe maak je een evaluatierapport?"
    ])

    add_content_slide(prs, "Waarom voortgang bewaken?", [
        "Zicht op of je op schema ligt",
        "Vroegtijdig problemen signaleren",
        "Keuzes maken bij vertraging",
        "Stakeholders informeren"
    ])

    add_content_slide(prs, "Voortgang monitoren", [
        "Welke user stories zijn afgerond?",
        "Welke zijn bezig?",
        "Welke moeten nog beginnen?",
        "Noteer de status van elke user story"
    ])

    add_content_slide(prs, "Planning controleren", [
        "Vergelijk werkelijke voortgang met planning",
        "Op schema, voor of achter?",
        "Noteer wat je ziet"
    ])

    add_content_slide(prs, "Prioriteiten bewaren", [
        "Must have eerst afgerond?",
        "Geen lagere prioriteiten als Must have nog open staat",
        "Documenteer of prioriteiten kloppen"
    ])

    add_content_slide(prs, "Afwijkingen documenteren", [
        "Wat is er anders dan gepland?",
        "Waarom is de afwijking ontstaan?",
        "Wat is de impact op de planning?",
        "Wat ga je doen om het op te lossen?"
    ])

    add_content_slide(prs, "Evaluatierapport", [
        "Inleiding: project en doel",
        "Voortgang: status van het project",
        "Planning: verloopt alles volgens plan?",
        "Prioriteiten: goed bewaakt?",
        "Afwijkingen en keuzes",
        "Conclusie: wat heb je geleerd?"
    ])

    add_content_slide(prs, "Opdracht deze sessie", [
        "Bewaak de voortgang van je project",
        "Documenteer afwijkingen",
        "Maak een evaluatierapport (2-3 pagina's)"
    ])

    add_content_slide(prs, "Online cursus", [
        "Agile Project Management with Jira",
        "LinkedIn Learning",
        "Projecten plannen en voortgang bewaken met Jira"
    ])

    add_title_slide(prs, "Afsluiting module", "Portfolio examen met echte klant – succes!")

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_DIR / "Sessie_5_Voortgang_bewaken_en_evalueren.pptx")


def main():
    print("Genereren van workshop presentaties...")
    create_workshop_1()
    print("  [OK] Sessie 1: Kennismaking en wensen/eisen verzamelen")
    create_workshop_2()
    print("  [OK] Sessie 2: User stories en backlog maken")
    create_workshop_3()
    print("  [OK] Sessie 3: Planning maken met MoSCoW")
    create_workshop_4()
    print("  [OK] Sessie 4: Ontwerpen met schema's")
    create_workshop_5()
    print("  [OK] Sessie 5: Voortgang bewaken en evalueren")
    print(f"\nAlle presentaties opgeslagen in: {OUTPUT_DIR}")


if __name__ == "__main__":
    main()
